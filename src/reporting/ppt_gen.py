from pptx import Presentation
from pptx.util import Inches
import os

class PPTReportBuilder:
    def __init__(self, template_path):
        if os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            print(f"Template not found: {template_path}. Creating new presentation.")
            self.prs = Presentation()
            
    def create_title_slide(self, title_text, subtitle_text):
        layout = self.prs.slide_layouts[0] # Title Slide
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle_text
            
    def create_summary_slide(self, df):
        layout = self.prs.slide_layouts[5] # Title Only
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = "Summary of Key Parameters"
        
        # Add table (simplified)
        cols_to_show = ['filename', 'J0', 'n', 'Rs', 'Rsh']
        cols_to_show = [c for c in cols_to_show if c in df.columns]
        
        rows = min(len(df), 10) # Limit rows
        cols = len(cols_to_show)
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table
        
        # Header
        for i, col in enumerate(cols_to_show):
            table.cell(0, i).text = str(col)
            
        # Rows
        for r in range(rows):
            row_data = df.iloc[r]
            for c, col_name in enumerate(cols_to_show):
                val = row_data[col_name]
                if isinstance(val, float):
                    txt = f"{val:.2e}" if abs(val) < 1e-3 or abs(val) > 1e3 else f"{val:.4f}"
                else:
                    txt = str(val)
                table.cell(r+1, c).text = txt

    def create_chart_slide(self, title, plot_stream, comments):
        layout = self.prs.slide_layouts[5] # Title Only
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        
        # Add Image
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(5.0)
        if plot_stream:
            slide.shapes.add_picture(plot_stream, left, top, height=height)
        
        # Add Comments
        txBox = slide.shapes.add_textbox(Inches(6), Inches(1.5), Inches(3.5), Inches(5))
        tf = txBox.text_frame
        tf.text = comments

    def save(self, filepath):
        self.prs.save(filepath)
