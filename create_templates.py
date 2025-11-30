from docx import Document
from pptx import Presentation
import os

if not os.path.exists('templates'):
    os.makedirs('templates')

# Word Template
doc = Document()
doc.add_heading('Dark IV Analysis Report Template', 0)
doc.save('templates/report_template.docx')

# PPT Template
prs = Presentation()
# Title slide layout
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Dark IV Analysis"
subtitle.text = "Template"
prs.save('templates/slide_template.pptx')

print("Templates created successfully.")
